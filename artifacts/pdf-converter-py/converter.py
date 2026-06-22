"""
========================================================
PDF CONVERTER TO DOCX, PPTX, XLSX
========================================================
Converts PDF pages to editable documents with OCR results
"""

import os
import uuid
import logging
from typing import List, Dict, Any
from pdf2image import convert_from_path
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pptx.enum.text import PP_ALIGN
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment

from ocr_engine import run_advanced_ocr, group_text_by_lines, analyze_layout

logger = logging.getLogger(__name__)

from PIL import Image
import subprocess
import shutil


# ========================================================
# PDF TO DOCX
# ========================================================

def convert_pdf_to_docx(pdf_path: str, output_path: str = None) -> str:
    """
    Convert PDF to DOCX with OCR and layout preservation
    
    Args:
        pdf_path: Path to PDF file
        output_path: Optional output path (default: UUID.docx)
    
    Returns:
        Path to generated DOCX file
    """
    try:
        file_id = output_path or f"{uuid.uuid4()}.docx"
        logger.info(f"Converting PDF to DOCX: {pdf_path} -> {file_id}")
        
        # Convert PDF to images
        logger.info("Converting PDF pages to images...")
        pages = convert_from_path(pdf_path, dpi=300)
        
        # Create document
        doc = Document()
        
        for page_num, page in enumerate(pages, 1):
            logger.info(f"Processing page {page_num}/{len(pages)}")
            
            # Save page as image
            image_path = f"temp_page_{page_num}.png"
            page.save(image_path, "PNG")
            
            try:
                # Run OCR
                extracted = run_advanced_ocr(image_path)
                
                if not extracted:
                    logger.warning(f"No text extracted from page {page_num}")
                    continue
                
                # Add page heading
                heading = doc.add_heading(f"Page {page_num}", level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                # Group text into lines
                lines = group_text_by_lines(extracted)
                
                # Add content
                for line in lines:
                    paragraph = doc.add_paragraph()
                    for item in line:
                        run = paragraph.add_run(item["text"] + " ")
                        run.font.size = Pt(11)
                        run.font.name = 'Calibri'
                        
                        # Add metadata comment if confidence is low
                        if item["confidence"] < 0.75:
                            logger.debug(f"Low confidence text: {item['text']} ({item['confidence']:.2%})")
                
                # Add page break
                doc.add_page_break()
            
            finally:
                # Cleanup
                if os.path.exists(image_path):
                    os.remove(image_path)
        
        # Save document
        doc.save(file_id)
        logger.info(f"DOCX saved to: {file_id}")
        
        return file_id
    
    except Exception as e:
        logger.error(f"PDF to DOCX conversion failed: {e}")
        raise


# ========================================================
# PDF TO PPTX
# ========================================================

def convert_pdf_to_pptx(pdf_path: str, output_path: str = None) -> str:
    """
    Convert PDF to PPTX with OCR and one slide per page
    
    Args:
        pdf_path: Path to PDF file
        output_path: Optional output path (default: UUID.pptx)
    
    Returns:
        Path to generated PPTX file
    """
    try:
        file_id = output_path or f"{uuid.uuid4()}.pptx"
        logger.info(f"Converting PDF to PPTX: {pdf_path} -> {file_id}")
        
        # Convert PDF to images
        logger.info("Converting PDF pages to images...")
        pages = convert_from_path(pdf_path, dpi=300)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = PPTXInches(10)
        prs.slide_height = PPTXInches(7.5)
        
        for page_num, page in enumerate(pages, 1):
            logger.info(f"Processing page {page_num}/{len(pages)}")
            
            # Save page as image
            image_path = f"temp_slide_{page_num}.png"
            page.save(image_path, "PNG")
            
            try:
                # Run OCR
                extracted = run_advanced_ocr(image_path)
                
                if not extracted:
                    logger.warning(f"No text extracted from page {page_num}")
                    continue
                
                # Add blank slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_box = slide.shapes.add_textbox(
                    PPTXInches(0.5),
                    PPTXInches(0.3),
                    PPTXInches(9),
                    PPTXInches(0.5)
                )
                title_frame = title_box.text_frame
                title_frame.text = f"Page {page_num}"
                title_frame.paragraphs[0].font.size = PPTXPt(28)
                title_frame.paragraphs[0].font.bold = True
                
                # Group text into lines
                lines = group_text_by_lines(extracted)
                
                # Add content
                top = PPTXInches(1.0)
                for line in lines:
                    text = " ".join([item["text"] for item in line])
                    
                    textbox = slide.shapes.add_textbox(
                        PPTXInches(0.5),
                        top,
                        PPTXInches(9),
                        PPTXInches(0.4)
                    )
                    
                    text_frame = textbox.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True
                    text_frame.paragraphs[0].font.size = PPTXPt(14)
                    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    
                    top += PPTXInches(0.5)
            
            finally:
                # Cleanup
                if os.path.exists(image_path):
                    os.remove(image_path)
        
        # Save presentation
        prs.save(file_id)
        logger.info(f"PPTX saved to: {file_id}")
        
        return file_id
    
    except Exception as e:
        logger.error(f"PDF to PPTX conversion failed: {e}")
        raise


# ========================================================
# PDF TO XLSX
# ========================================================

def convert_pdf_to_xlsx(pdf_path: str, output_path: str = None) -> str:
    """
    Convert PDF to XLSX with OCR data structured in cells
    
    Args:
        pdf_path: Path to PDF file
        output_path: Optional output path (default: UUID.xlsx)
    
    Returns:
        Path to generated XLSX file
    """
    try:
        file_id = output_path or f"{uuid.uuid4()}.xlsx"
        logger.info(f"Converting PDF to XLSX: {pdf_path} -> {file_id}")
        
        # Convert PDF to images
        logger.info("Converting PDF pages to images...")
        pages = convert_from_path(pdf_path, dpi=300)
        
        # Create workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "OCR Data"
        
        # Add headers
        headers = ["Page", "Line #", "Text", "Confidence", "Language", "X Position", "Y Position"]
        ws.append(headers)
        
        # Style headers
        for cell in ws[1]:
            cell.font = Font(bold=True, size=12)
            cell.alignment = Alignment(horizontal='center')
        
        row = 2
        
        for page_num, page in enumerate(pages, 1):
            logger.info(f"Processing page {page_num}/{len(pages)}")
            
            # Save page as image
            image_path = f"temp_excel_{page_num}.png"
            page.save(image_path, "PNG")
            
            try:
                # Run OCR
                extracted = run_advanced_ocr(image_path)
                
                if not extracted:
                    logger.warning(f"No text extracted from page {page_num}")
                    continue
                
                # Group text into lines
                lines = group_text_by_lines(extracted)
                
                # Add data
                for line_num, line in enumerate(lines, 1):
                    text = " ".join([item["text"] for item in line])
                    confidence = sum([item["confidence"] for item in line]) / len(line)
                    language = line[0]["language"] if line else "unknown"
                    x_pos = line[0]["x"] if line else 0
                    y_pos = line[0]["y"] if line else 0
                    
                    ws.append([
                        page_num,
                        line_num,
                        text,
                        f"{confidence:.2%}",
                        language,
                        f"{x_pos:.0f}",
                        f"{y_pos:.0f}"
                    ])
                    
                    row += 1
            
            finally:
                # Cleanup
                if os.path.exists(image_path):
                    os.remove(image_path)
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 8
        ws.column_dimensions['B'].width = 10
        ws.column_dimensions['C'].width = 40
        ws.column_dimensions['D'].width = 12
        ws.column_dimensions['E'].width = 12
        ws.column_dimensions['F'].width = 12
        ws.column_dimensions['G'].width = 12
        
        # Save workbook
        wb.save(file_id)
        logger.info(f"XLSX saved to: {file_id}")
        
        return file_id
    
    except Exception as e:
        logger.error(f"PDF to XLSX conversion failed: {e}")
        raise


# ========================================================
# MAIN CONVERTER
# ========================================================

def convert_pdf(pdf_path: str, output_format: str, output_path: str = None) -> str:
    """
    Main converter function
    
    Args:
        pdf_path: Path to PDF file
        output_format: Output format ('docx', 'pptx', or 'xlsx')
        output_path: Optional output path
    
    Returns:
        Path to converted file
    """
    logger.info(f"Starting PDF conversion: {pdf_path} -> {output_format}")
    
    if output_format.lower() == "docx":
        return convert_pdf_to_docx(pdf_path, output_path)
    
    elif output_format.lower() == "pptx":
        return convert_pdf_to_pptx(pdf_path, output_path)
    
    elif output_format.lower() == "xlsx":
        return convert_pdf_to_xlsx(pdf_path, output_path)
    
    else:
        raise ValueError(f"Unsupported format: {output_format}")


# ========================================================
# ANY/IMAGE/DOC -> PDF
# ========================================================

def convert_to_pdf(input_path: str, output_path: str = None) -> str:
    """
    Convert various input documents (images, office files) to PDF.

    - Images (png, jpg, jpeg, tiff, bmp, gif) are converted with Pillow.
    - Office documents (docx, pptx, xlsx, odt, ods, odp, etc.) are converted
      using LibreOffice `soffice` if available.

    Args:
        input_path: Path to input file
        output_path: Optional output PDF path

    Returns:
        Path to generated PDF file
    """
    try:
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input not found: {input_path}")

        ext = os.path.splitext(input_path)[1].lower()
        image_exts = {'.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.gif'}

        out_file = output_path or f"{uuid.uuid4()}.pdf"

        # IMAGE -> PDF (supports multi-frame TIFF)
        if ext in image_exts:
            logger.info(f"Converting image to PDF: {input_path} -> {out_file}")
            img = Image.open(input_path)

            # Convert RGBA to RGB for PDF
            def _ensure_rgb(im):
                if im.mode in ("RGBA", "LA"):
                    bg = Image.new("RGB", im.size, (255, 255, 255))
                    bg.paste(im, mask=im.split()[-1])
                    return bg
                return im.convert("RGB")

            if getattr(img, "is_animated", False) or getattr(img, "n_frames", 1) > 1:
                frames = []
                for i in range(getattr(img, "n_frames", 1)):
                    img.seek(i)
                    frames.append(_ensure_rgb(img.copy()))

                frames[0].save(out_file, "PDF", save_all=True, append_images=frames[1:])
            else:
                rgb = _ensure_rgb(img)
                rgb.save(out_file, "PDF")

            logger.info(f"PDF saved to: {out_file}")
            return out_file

        # OFFICE/OTHER -> PDF via LibreOffice
        else:
            logger.info(f"Attempting to convert document to PDF using LibreOffice: {input_path}")
            soffice = shutil.which("soffice")
            if not soffice:
                raise RuntimeError(
                    "LibreOffice `soffice` not found. Install LibreOffice to enable document-to-PDF conversion."
                )

            # LibreOffice writes output into the specified outdir
            outdir = os.path.abspath(os.path.dirname(out_file))
            os.makedirs(outdir, exist_ok=True)

            cmd = [soffice, "--headless", "--convert-to", "pdf", input_path, "--outdir", outdir]
            logger.info(f"Running: {' '.join(cmd)}")
            res = subprocess.run(cmd, capture_output=True, text=True)

            if res.returncode != 0:
                logger.error(res.stdout)
                logger.error(res.stderr)
                raise RuntimeError(f"LibreOffice conversion failed: {res.stderr.strip()}")

            # LibreOffice will create a file with same stem but .pdf
            generated = os.path.join(outdir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")

            if not os.path.exists(generated):
                raise RuntimeError("LibreOffice did not produce a PDF output as expected")

            # Move/rename to requested output_path if necessary
            if os.path.abspath(generated) != os.path.abspath(out_file):
                os.replace(generated, out_file)

            logger.info(f"PDF saved to: {out_file}")
            return out_file

    except Exception as e:
        logger.error(f"Conversion to PDF failed: {e}")
        raise
